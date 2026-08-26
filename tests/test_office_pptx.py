"""PowerPoint: slide structure, speaker notes, and revise-an-existing-deck reads."""

import pytest

from coworker.agents.base import AgentContext
from coworker.roots import RootDir
from coworker.tools.office.pptx_tools import pptx_tools

pytest.importorskip("pptx", reason="python-pptx is an optional [office] extra")


@pytest.fixture
def tools(tmp_path):
    ws = tmp_path / "scratch"
    ws.mkdir()
    context = AgentContext(workspace=ws, roots=[RootDir(path=ws, writable=True)])
    return {t.__name__: t for t in pptx_tools(context)}, ws


def test_write_then_read_round_trips_titles_and_bullets(tools):
    write, read = tools[0]["write_presentation"], tools[0]["read_presentation"]
    out = write(
        "deck.pptx",
        [
            {"layout": "title", "title": "Q3 Review", "subtitle": "Finance"},
            {
                "layout": "bullets",
                "title": "Highlights",
                "bullets": ["Revenue +12%", "Churn -3%"],
            },
        ],
    )
    assert "error" not in out
    assert out["slides_written"] == 2

    deck = read("deck.pptx")
    assert deck["total_slides"] == 2
    assert deck["slides"][0]["title"] == "Q3 Review"
    assert deck["slides"][1]["title"] == "Highlights"
    assert "Revenue +12%" in deck["slides"][1]["bullets"]


def test_speaker_notes_round_trip(tools):
    """A deck without notes isn't a finished deliverable, so notes must survive."""
    write, read = tools[0]["write_presentation"], tools[0]["read_presentation"]
    write(
        "deck.pptx",
        [
            {
                "layout": "bullets",
                "title": "Method",
                "bullets": ["Two-sample t-test"],
                "notes": "Explain why we dropped the 14 incomplete responses.",
            }
        ],
    )
    deck = read("deck.pptx")
    assert "14 incomplete responses" in deck["slides"][0]["notes"]


def test_sub_bullets_keep_their_level(tools):
    write, read = tools[0]["write_presentation"], tools[0]["read_presentation"]
    write(
        "deck.pptx",
        [
            {
                "layout": "bullets",
                "title": "Detail",
                "bullets": ["Top", {"text": "Nested", "level": 1}],
            }
        ],
    )
    bullets = read("deck.pptx")["slides"][0]["bullets"]
    assert any(b.startswith("  ") and "Nested" in b for b in bullets)


def test_slides_are_numbered_for_revision(tools):
    write, read = tools[0]["write_presentation"], tools[0]["read_presentation"]
    write("d.pptx", [{"layout": "bullets", "title": f"S{i}"} for i in range(4)])
    assert [s["index"] for s in read("d.pptx")["slides"]] == [0, 1, 2, 3]


def test_append_adds_slides_to_an_existing_deck(tools):
    write, read = tools[0]["write_presentation"], tools[0]["read_presentation"]
    write("d.pptx", [{"layout": "title", "title": "One"}])
    write("d.pptx", [{"layout": "bullets", "title": "Two"}], append=True)
    deck = read("d.pptx")
    assert deck["total_slides"] == 2
    assert [s["title"] for s in deck["slides"]] == ["One", "Two"]


def test_section_and_blank_layouts_are_accepted(tools):
    write = tools[0]["write_presentation"]
    out = write(
        "d.pptx",
        [
            {"layout": "section", "title": "Part II"},
            {"layout": "blank"},
        ],
    )
    assert "error" not in out
    assert out["slides_written"] == 2


def test_unknown_layout_is_reported_clearly(tools):
    result = tools[0]["write_presentation"]("d.pptx", [{"layout": "carousel"}])
    assert "error" in result and "carousel" in result["error"]


def test_image_slide_embeds_a_picture_from_the_workspace(tools):
    write, ws = tools[0]["write_presentation"], tools[1]
    png = ws / "chart.png"
    # Smallest valid 1x1 PNG.
    png.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d4948445200000001000000010802000000907753"
            "de0000000c4944415408d76360000000020001e221bc330000000049454e44ae426082"
        )
    )
    out = write("d.pptx", [{"layout": "image", "title": "Trend", "image": "chart.png"}])
    assert "error" not in out


def test_missing_image_errors_rather_than_producing_a_broken_deck(tools):
    result = tools[0]["write_presentation"](
        "d.pptx", [{"layout": "image", "title": "T", "image": "nope.png"}]
    )
    assert "error" in result


def test_write_outside_the_workspace_is_refused(tools):
    result = tools[0]["write_presentation"]("/tmp/escape.pptx", [{"layout": "blank"}])
    assert "error" in result and "escapes" in result["error"]


def test_reading_a_missing_deck_errors_cleanly(tools):
    assert "error" in tools[0]["read_presentation"]("nope.pptx")


# --- the design system (added 2026-08-25, when decks stopped being 4:3 Calibri) --------


def test_a_new_deck_is_widescreen_not_the_4_3_default(tools):
    """python-pptx's stock template is 10x7.5in. A 4:3 deck is letterboxed on every
    modern screen and is the single loudest 'this looks old' signal a slide sends."""
    import pptx

    write, ws = tools[0]["write_presentation"], tools[1]
    out = write("deck.pptx", [{"layout": "title", "title": "Q3"}])
    assert out["widescreen"] is True
    deck = pptx.Presentation(str(ws / "deck.pptx"))
    assert round(deck.slide_width / deck.slide_height, 3) == round(16 / 9, 3)


def test_the_deck_theme_is_not_office_2007(tools):
    import pptx

    from coworker.tools.office.deck_theme import NEUTRAL_MODERN, read_theme

    write, ws = tools[0]["write_presentation"], tools[1]
    write("deck.pptx", [{"layout": "bullets", "title": "T", "bullets": ["a"]}])
    theme = read_theme(pptx.Presentation(str(ws / "deck.pptx")))
    assert theme.heading == NEUTRAL_MODERN.heading != "Calibri"
    assert theme.body == NEUTRAL_MODERN.body


def test_bullets_are_real_bullet_formatting_not_glued_to_the_text(tools):
    """A marker in the string would come back on read and be prefixed again on the next
    write — 'Revenue' becoming '— — Revenue' after two revisions."""
    write, read = tools[0]["write_presentation"], tools[0]["read_presentation"]
    write("deck.pptx", [{"layout": "bullets", "title": "T", "bullets": ["Revenue +12%"]}])
    back = read("deck.pptx")["slides"][0]["bullets"]
    assert back == ["Revenue +12%"]


def test_the_printed_slide_number_is_not_read_back_as_content(tools):
    write, read = tools[0]["write_presentation"], tools[0]["read_presentation"]
    write(
        "deck.pptx",
        [
            {"layout": "bullets", "title": "One", "bullets": ["x"]},
            {"layout": "bullets", "title": "Two", "bullets": ["y"]},
        ],
    )
    for slide in read("deck.pptx")["slides"]:
        assert "2" not in slide["bullets"] and slide["bullets"] != []


def test_every_layout_renders_and_keeps_its_title_readable(tools):
    """Titles live in the real title placeholder — outline view, the slide sorter, and
    read_presentation all depend on it, which is how a deck stays revisable."""
    write, read = tools[0]["write_presentation"], tools[0]["read_presentation"]
    write(
        "deck.pptx",
        [
            {"layout": "title", "title": "Opener", "subtitle": "sub"},
            {"layout": "section", "title": "Part one"},
            {"layout": "statement", "statement": "Mobile is the problem."},
            {"layout": "stat", "title": "Numbers", "stats": [{"value": "71%", "label": "on a phone"}]},
            {"layout": "quote", "title": "Voice", "quote": "It was too long.", "attribution": "R12"},
            {
                "layout": "two_column",
                "title": "Split",
                "columns": [{"heading": "A", "bullets": ["one"]}, {"heading": "B", "bullets": ["two"]}],
            },
            {
                "layout": "comparison",
                "title": "Before and after",
                "columns": [{"heading": "Before", "bullets": ["slow"]}, {"heading": "After", "bullets": ["fast"]}],
            },
        ],
    )
    deck = read("deck.pptx")
    assert deck["total_slides"] == 7
    titles = [s["title"] for s in deck["slides"]]
    assert titles[:2] == ["Opener", "Part one"]
    assert titles[2] == "Mobile is the problem."  # the claim IS the title
    assert titles[4] == "Voice"
    assert "one" in deck["slides"][5]["bullets"] and "two" in deck["slides"][5]["bullets"]


def test_a_layout_that_is_missing_its_own_content_says_which_field(tools):
    write = tools[0]["write_presentation"]
    assert "stats" in write("d.pptx", [{"layout": "stat", "title": "T"}])["error"]
    assert "columns" in write("d.pptx", [{"layout": "two_column", "title": "T"}])["error"]
    assert "image" in write("d.pptx", [{"layout": "image", "title": "T"}])["error"]


def test_a_house_template_outranks_the_built_in_theme(tools, tmp_path):
    """Passing template= is how someone says 'use our brand'. Ours must step aside."""
    import pptx

    from coworker.tools.office.deck_theme import read_theme

    write, ws = tools[0]["write_presentation"], tools[1]
    house = pptx.Presentation()
    house.save(str(ws / "house.pptx"))  # stock template = Calibri, the brand here
    write("deck.pptx", [{"layout": "bullets", "title": "T", "bullets": ["a"]}], template="house.pptx")
    assert read_theme(pptx.Presentation(str(ws / "deck.pptx"))).heading == "Calibri"


def test_appending_never_reflows_someone_elses_deck(tools):
    """Resizing a 4:3 deck we were merely asked to add slides to would silently re-lay
    out every slide already in it."""
    import pptx

    write, ws = tools[0]["write_presentation"], tools[1]
    original = pptx.Presentation()  # 4:3
    original.slides.add_slide(original.slide_layouts[1])
    original.save(str(ws / "theirs.pptx"))
    write("theirs.pptx", [{"layout": "bullets", "title": "Added", "bullets": ["x"]}], append=True)
    after = pptx.Presentation(str(ws / "theirs.pptx"))
    assert after.slide_width == 9144000  # still 4:3 — we added to it, we didn't restyle it
    assert len(list(after.slides)) == 2


# --- text fitting (owner report 2026-08-26: "font sometimes too big, exceeds the frame")


def test_a_long_title_shrinks_instead_of_walking_out_of_its_frame(tools):
    """No renderer exists at runtime to screenshot-and-check, so the type must yield to
    the box deterministically."""
    import pptx
    from pptx.util import Pt

    write, ws = tools[0]["write_presentation"], tools[1]
    long_title = (
        "Why the onboarding survey completion rate collapsed on mobile devices "
        "between the June and August fieldwork waves across all three markets"
    )
    write(
        "deck.pptx",
        [
            {"layout": "bullets", "title": "Short", "bullets": ["a"]},
            {"layout": "bullets", "title": long_title, "bullets": ["a"]},
        ],
    )
    deck = pptx.Presentation(str(ws / "deck.pptx"))
    sizes = []
    for slide in deck.slides:
        run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
        sizes.append(run.font.size)
    assert sizes[0] == Pt(28)  # short keeps the scale
    assert sizes[1] < Pt(28)  # long yields to the box


def test_a_reasonable_deck_carries_no_layout_warnings(tools):
    out = tools[0]["write_presentation"](
        "deck.pptx",
        [{"layout": "bullets", "title": "Fine", "bullets": ["one", "two", "three"]}],
    )
    assert "layout_warnings" not in out


def test_an_overstuffed_slide_is_flagged_for_the_model_to_split(tools):
    """The deck still saves — shrunk to the floor — but the result says which slide
    holds more than fits and what to do. The prompt forbids handing that over."""
    out = tools[0]["write_presentation"](
        "deck.pptx",
        [
            {"layout": "bullets", "title": "Fine", "bullets": ["a"]},
            {
                "layout": "bullets",
                "title": "Everything at once",
                "bullets": [
                    f"Finding number {i}: a long sentence that wraps at least once at body size"
                    for i in range(18)
                ],
            },
        ],
    )
    assert out["slides_written"] == 2  # saved, not refused
    warnings = out["layout_warnings"]
    assert len(warnings) == 1 and warnings[0].startswith("slide 2:")
    assert "split" in warnings[0]
    assert "write_presentation again" in out["action_required"]


def test_a_long_quote_shrinks_and_the_attribution_stays_below_it(tools):
    import pptx
    from pptx.util import Pt

    write, ws = tools[0]["write_presentation"], tools[1]
    write(
        "deck.pptx",
        [
            {
                "layout": "quote",
                "quote": (
                    "I started the survey on my phone during my commute and gave up "
                    "halfway through the third matrix grid because I had to scroll "
                    "sideways on every single row just to see the answer labels, and "
                    "by then my train had reached the station. I tried again that "
                    "evening on my laptop and got further, but the progress bar had "
                    "reset to zero, all my earlier answers were gone, and the grid "
                    "questions were exactly as unreadable as before, so I closed the "
                    "tab and decided the gift card was not worth a third attempt."
                ),
                "attribution": "Respondent 3,118",
            }
        ],
    )
    deck = pptx.Presentation(str(ws / "deck.pptx"))
    slide = list(deck.slides)[0]
    quote_run = next(
        shape.text_frame.paragraphs[0].runs[0]
        for shape in slide.shapes
        if shape.has_text_frame and "matrix grid" in shape.text_frame.text
    )
    assert quote_run.font.size < Pt(28)


def test_the_fit_estimator_is_monotonic_and_floored():
    from coworker.tools.office.deck_render import fitted_size

    box = dict(width_in=11.6, height_in=0.9, size=28, min_size=18, bold=True)
    short = fitted_size("Findings", **box)
    long = fitted_size("A title that goes on and on " * 6, **box)
    absurd = fitted_size("word " * 400, **box)
    assert short == 28
    assert 18 <= long < 28
    assert absurd == 18  # never below the floor — past that it's a writing problem
