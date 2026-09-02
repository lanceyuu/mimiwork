"""Word and PowerPoint read back as a preview you can point at (owner ask 2026-09-02)."""

from __future__ import annotations

import struct
import zlib

import pytest

from coworker.office_preview import add_word_comment, docx_to_html, pptx_to_html


def _png_1x1() -> bytes:
    def chunk(tag, data):
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

    raw = zlib.compress(b"\x00\xff\x00\x00")
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)) + chunk(b"IDAT", raw) + chunk(b"IEND", b"")


@pytest.fixture()
def docx_path(tmp_path):
    from docx import Document

    doc = Document()
    doc.add_heading("Findings", level=1)
    p = doc.add_paragraph("The sample ")
    p.add_run("consisted").bold = True
    p.add_run(" of 24 interviews & 2 groups.")
    doc.add_paragraph("First theme", style="List Bullet")
    doc.add_paragraph("Second theme", style="List Bullet")
    t = doc.add_table(rows=2, cols=2)
    t.cell(0, 0).text, t.cell(0, 1).text = "Code", "Count"
    t.cell(1, 0).text, t.cell(1, 1).text = "Trust", "12"
    png = tmp_path / "dot.png"
    png.write_bytes(_png_1x1())
    doc.add_picture(str(png))
    path = tmp_path / "report.docx"
    doc.save(str(path))
    return path


def test_a_word_file_reads_back_in_order_with_paragraph_indexes(docx_path):
    r = docx_to_html(docx_path)
    h = r["html"]
    assert '<h1 data-p="0">Findings</h1>' in h
    assert "<p data-p=\"1\">The sample <b>consisted</b> of 24 interviews &amp; 2 groups.</p>" in h
    assert "<ul>" in h and '<li data-p="2">First theme</li>' in h and '<li data-p="3">Second theme</li>' in h
    assert "<table><tr><th>Code</th><th>Count</th></tr><tr><td>Trust</td><td>12</td></tr></table>" in h
    assert '<img src="data:image/png;base64,' in h
    assert r["paragraphs"] >= 5


def test_a_comment_lands_on_the_paragraph_as_a_real_word_comment(docx_path):
    from docx import Document

    info = add_word_comment(docx_path, 1, "Say how they were recruited.", author="Shubin Yu")
    assert info["comments"] == 1
    doc = Document(str(docx_path))
    comments = list(doc.comments)
    assert comments[0].text == "Say how they were recruited."
    assert comments[0].author == "Shubin Yu" and comments[0].initials == "SY"
    with pytest.raises(ValueError):
        add_word_comment(docx_path, 99, "nowhere")


def test_slides_read_back_one_card_each_with_notes(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Method"
    s.placeholders[1].text_frame.text = "24 interviews"
    s.placeholders[1].text_frame.add_paragraph().text = "2 focus groups"
    s.notes_slide.notes_text_frame.text = "Mention the pilot."
    prs.slides.add_slide(prs.slide_layouts[5]).shapes.title.text = "Results"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    r = pptx_to_html(path)
    assert r["slides"] == 2
    assert '<section class="slide" data-slide="1">' in r["html"]
    assert "<h2>Method</h2>" in r["html"] and "<p>24 interviews</p>" in r["html"]
    assert "Notes: Mention the pilot." in r["html"]
    assert 'data-slide="2"' in r["html"] and "<h2>Results</h2>" in r["html"]


def test_the_session_reader_serves_them_and_refuses_comments_elsewhere(tmp_path, monkeypatch, docx_path):
    from helpers import CapturingProvider

    from coworker.server.manager import SessionManager

    monkeypatch.setenv("COWORKER_STATE_DIR", str(tmp_path / "state"))
    m = SessionManager(workspace=tmp_path, data_dir=tmp_path / "data", provider=CapturingProvider())
    engine = m.get_engine("s1", agent="cowork")
    assert engine is not None
    got = m.read_artifact("s1", docx_path.name)
    assert got["kind"] == "docx" and "Findings" in got["content"] and got["paragraphs"] >= 5
    assert m.comment_artifact("s1", docx_path.name, 0, "Title needs a date", author="")["ok"]
    assert not m.comment_artifact("s1", docx_path.name, 0, "   ")["ok"]
    assert not m.comment_artifact("s1", "nope.txt", 0, "x")["ok"]
