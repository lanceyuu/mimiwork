"""PowerPoint decks — read and write .pptx deliverables.

Slides are written from a small IR (``title`` / ``bullets`` / ``section`` / ``image`` /
``blank``) for the same reason Word is: the model produces structured JSON reliably and
OOXML not at all.

Speaker notes are a first-class field, not an afterthought. A deck handed over without notes
is not a finished deliverable — the person presenting it has to reconstruct the argument from
bullet fragments.

``read_presentation`` returns per-slide title, bullets, and notes so an existing deck can be
revised rather than regenerated from scratch.
"""

from __future__ import annotations

from typing import Any

from ... import deliverable_check
from ._common import clip, decorate, guard, require
from .paths import context_roots, display_path, resolve_read, resolve_write

# python-pptx's default template layout indexes.
_LAYOUT_TITLE = 0
_LAYOUT_TITLE_CONTENT = 1
_LAYOUT_SECTION = 2
_LAYOUT_TITLE_ONLY = 5
_LAYOUT_BLANK = 6

_SLIDE_SHAPE = {
    "type": "object",
    "properties": {
        "layout": {
            "type": "string",
            "enum": ["title", "bullets", "section", "image", "blank"],
            "description": (
                "Slide kind: 'title' (deck opener), 'bullets' (title + bullet body), "
                "'section' (divider), 'image' (title + picture), 'blank'."
            ),
        },
        "title": {"type": "string", "description": "Slide title."},
        "subtitle": {"type": "string", "description": "Subtitle (title/section layouts)."},
        "bullets": {
            "type": "array",
            "description": "Body bullets. Use {'text': ..., 'level': 1} for sub-bullets.",
            "items": {},
        },
        "image": {
            "type": "string",
            "description": "Path to a PNG/JPG to place on the slide (image layout).",
        },
        "notes": {
            "type": "string",
            "description": "Speaker notes — what the presenter should say on this slide.",
        },
    },
}

_WRITE_SCHEMA = {
    "type": "function",
    "function": {
        "name": "write_presentation",
        "description": (
            "Create or overwrite a PowerPoint (.pptx) deck from structured slides. Always "
            "include speaker notes — a deck without notes is not a finished deliverable. Use "
            "this for any slide deliverable; do NOT write a script to do it. Pass append=true "
            "to add slides to an existing deck."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Destination .pptx path."},
                "slides": {
                    "type": "array",
                    "description": "Slides, in order.",
                    "items": _SLIDE_SHAPE,
                },
                "append": {
                    "type": "boolean",
                    "description": "Append to an existing deck (default false = overwrite).",
                },
                "template": {
                    "type": "string",
                    "description": (
                        "Optional .pptx/.potx to inherit theme, fonts, and branding from. "
                        "Use the organisation's template when there is one."
                    ),
                },
            },
            "required": ["path", "slides"],
        },
    },
}

_READ_SCHEMA = {
    "type": "function",
    "function": {
        "name": "read_presentation",
        "description": (
            "Read a PowerPoint (.pptx) deck as numbered slides with their titles, bullet text, "
            "and speaker notes — so an existing deck can be revised rather than rebuilt. "
            "Read-only."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "The .pptx file to read."},
            },
            "required": ["path"],
        },
    },
}


def _bullet_entry(item: Any) -> tuple[str, int]:
    if isinstance(item, dict):
        level = item.get("level", 0)
        return str(item.get("text") or ""), level if isinstance(level, int) and level > 0 else 0
    return str(item), 0


def _set_notes(slide: Any, notes: str) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def pptx_tools(context: Any) -> list:
    roots = context_roots(context)

    @guard
    def write_presentation(
        path: str, slides: list, append: bool = False, template: str = ""
    ) -> dict[str, Any]:
        pptx = require("pptx", "python-pptx")
        from pptx.util import Inches

        target = resolve_write(path, roots)
        if not isinstance(slides, list):
            raise ValueError("'slides' must be a list of slide objects")

        if append and target.is_file():
            deck = pptx.Presentation(str(target))
        elif template:
            deck = pptx.Presentation(str(resolve_read(template, roots)))
        else:
            deck = pptx.Presentation()

        for entry in slides:
            if not isinstance(entry, dict):
                raise ValueError(f"each slide must be an object, got {type(entry).__name__}")
            kind = str(entry.get("layout") or "bullets").lower()
            title = str(entry.get("title") or "")
            subtitle = str(entry.get("subtitle") or "")
            bullets = entry.get("bullets") or []

            if kind == "title":
                layout = _LAYOUT_TITLE
            elif kind == "section":
                layout = _LAYOUT_SECTION
            elif kind == "blank":
                layout = _LAYOUT_BLANK
            elif kind == "image":
                layout = _LAYOUT_TITLE_ONLY
            elif kind == "bullets":
                layout = _LAYOUT_TITLE_CONTENT
            else:
                raise ValueError(
                    f"unknown slide layout {kind!r}; expected one of: title, bullets, "
                    "section, image, blank"
                )
            # A user-supplied template may not carry all of python-pptx's default layouts.
            if layout >= len(deck.slide_layouts):
                layout = min(_LAYOUT_TITLE_CONTENT, len(deck.slide_layouts) - 1)

            slide = deck.slides.add_slide(deck.slide_layouts[layout])

            if title and slide.shapes.title is not None:
                slide.shapes.title.text = title

            if subtitle:
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == 1:
                        placeholder.text = subtitle
                        break

            if bullets and kind == "bullets":
                body = None
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx != 0:
                        body = placeholder
                        break
                if body is not None:
                    frame = body.text_frame
                    frame.clear()
                    for i, item in enumerate(bullets):
                        text, level = _bullet_entry(item)
                        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                        para.text = text
                        para.level = min(level, 4)

            image = entry.get("image")
            if image and kind == "image":
                picture = resolve_read(str(image), roots)
                if not picture.is_file():
                    raise FileNotFoundError(display_path(picture, roots))
                # Centred under the title, sized to the slide with a margin; python-pptx keeps
                # the aspect ratio when only a width is given.
                slide.shapes.add_picture(
                    str(picture), Inches(1.0), Inches(1.8), width=deck.slide_width - Inches(2.0)
                )

            _set_notes(slide, str(entry.get("notes") or ""))

        target.parent.mkdir(parents=True, exist_ok=True)
        deck.save(str(target))
        return deliverable_check.attach(
            {
                "path": display_path(target, roots),
                "slides_written": len(slides),
                "total_slides": len(list(deck.slides)),
                "appended": bool(append),
                "bytes": target.stat().st_size,
            },
            target,
        )

    @guard
    def read_presentation(path: str) -> dict[str, Any]:
        pptx = require("pptx", "python-pptx")
        target = resolve_read(path, roots)
        if not target.is_file():
            raise FileNotFoundError(display_path(target, roots))

        deck = pptx.Presentation(str(target))
        slides: list[dict[str, Any]] = []
        for index, slide in enumerate(deck.slides):
            title = ""
            if slide.shapes.title is not None:
                title = slide.shapes.title.text.strip()

            bullets: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame or shape == slide.shapes.title:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        bullets.append(clip(("  " * para.level) + text))

            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            slides.append(
                {
                    "index": index,
                    "title": clip(title),
                    "bullets": bullets,
                    "notes": clip(notes),
                }
            )

        return {
            "path": display_path(target, roots),
            "total_slides": len(slides),
            "slides": slides,
        }

    return [
        decorate(
            write_presentation,
            name="write_presentation",
            schema=_WRITE_SCHEMA,
            risk="medium",
            capabilities=["write"],
        ),
        decorate(read_presentation, name="read_presentation", schema=_READ_SCHEMA),
    ]
