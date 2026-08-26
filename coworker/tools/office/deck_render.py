"""Painting a slide: one function per slide kind, all sharing one design system.

Every layout here is drawn explicitly — position, size, face, weight, colour — rather than
poured into python-pptx's stock placeholders. That is the whole point: the placeholders are
where the 2007 look comes from. The title still uses a real title placeholder so PowerPoint's
outline view and slide sorter keep working; everything else is drawn.

The layout vocabulary exists because a deck made only of bullet lists is a deck nobody wants
to sit through. `statement`, `stat`, `quote`, `two_column` and `comparison` are the shapes an
argument actually takes — a claim, a number that carries it, someone's own words, a
side-by-side. Giving the model those shapes is what stops every slide becoming five bullets.
"""

from __future__ import annotations

from typing import Any, Callable, Optional

from .deck_theme import Theme

LAYOUTS = (
    "title",
    "section",
    "bullets",
    "statement",
    "stat",
    "two_column",
    "comparison",
    "quote",
    "image",
    "blank",
)


import math

# Average glyph width as a fraction of the point size, for the faces this module sets
# (Georgia/Arial-class). Deliberately a little wide: overestimating wraps a line early,
# underestimating lets text walk out of its frame — only one of those is visible.
_CHAR_FACTOR = 0.52
_CHAR_FACTOR_BOLD = 0.56


def text_lines(text: str, width_in: float, size: int, *, bold: bool = False) -> int:
    """Estimated wrapped-line count for `text` set at `size` in a `width_in` box."""
    factor = _CHAR_FACTOR_BOLD if bold else _CHAR_FACTOR
    per_line = max(1, int((width_in * 72.0) / (factor * size)))
    # Words wrap whole, so real text needs a bit more room than the char count says.
    return max(1, math.ceil(len(text) * 1.08 / per_line))


def block_height(
    text: str, width_in: float, size: int, *, bold: bool = False, line_spacing: float = 1.15
) -> float:
    """Estimated height (inches) of `text` wrapped in a `width_in` box at `size`."""
    return text_lines(text, width_in, size, bold=bold) * size * line_spacing / 72.0


def fitted_size(
    text: str,
    width_in: float,
    height_in: float,
    size: int,
    *,
    min_size: int,
    bold: bool = False,
    line_spacing: float = 1.15,
) -> int:
    """The largest size ≤ `size` (but ≥ `min_size`) whose text fits the box.

    This is the module's answer to "the font sometimes exceeds the frame": every box is
    fixed, so the type must yield. It cannot look at pixels — users' machines have no
    renderer to screenshot with — so it estimates conservatively and shrinks. Content
    that still overflows AT `min_size` is a writing problem, not a typesetting one;
    callers surface that as a layout warning so the model splits the slide.
    """
    while size > min_size and block_height(
        text, width_in, size, bold=bold, line_spacing=line_spacing
    ) > height_in:
        size -= 1
    return size


def _rgb(value: int):
    from pptx.dml.color import RGBColor

    return RGBColor((value >> 16) & 0xFF, (value >> 8) & 0xFF, value & 0xFF)


def _pick_layout(deck: Any, *names: str):
    """A layout by name, so a user's template resolves properly, with an index fallback."""
    wanted = {n.lower() for n in names}
    for layout in deck.slide_layouts:
        if (layout.name or "").strip().lower() in wanted:
            return layout
    index = 6 if "blank" in wanted else 5
    return deck.slide_layouts[min(index, len(deck.slide_layouts) - 1)]


def _blank_slide(deck: Any):
    return deck.slides.add_slide(_pick_layout(deck, "Blank"))


def _drop_empty_placeholders(slide: Any) -> None:
    """A placeholder left empty still prints 'Click to add text' in edit view."""
    for shape in list(slide.placeholders):
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            shape._element.getparent().remove(shape._element)


def _box(slide: Any, x: float, y: float, w: float, h: float):
    from pptx.util import Inches

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.word_wrap = True
    return frame


def _write(
    frame: Any,
    text: str,
    *,
    font: str,
    size: int,
    color: int,
    bold: bool = False,
    italic: bool = False,
    first: bool = False,
    space_after: int = 0,
    space_before: int = 0,
    line_spacing: Optional[float] = None,
    level: int = 0,
    align: Any = None,
):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.text = text
    # Explicit, always: a placeholder inherits its layout's alignment, and the stock
    # "Title Only" layout centres. Inheriting that is what made every title centred.
    para.alignment = align if align is not None else PP_ALIGN.LEFT
    para.level = min(max(level, 0), 4)
    if space_after:
        para.space_after = Pt(space_after)
    if space_before:
        para.space_before = Pt(space_before)
    if line_spacing:
        para.line_spacing = line_spacing
    for run in para.runs or [para.add_run()]:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
    return para


FOOTER_SHAPE_NAME = "mimiwork-slide-number"


def _bullet_marker(para: Any, char: str, face: str, color: int, indent: float) -> None:
    """Give a paragraph a real PowerPoint bullet.

    Prefixing the text with "—  " would look identical and be wrong: the marker becomes
    part of the string, so reading the deck back to revise it returns "—  Revenue +12%"
    and the next write prefixes it again. Bullets are paragraph formatting; treat them
    as such.
    """
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(indent))))
    pPr.set("indent", str(-int(Inches(indent))))
    for tag, attrs in (
        ("a:buClr", None),
        ("a:buFont", {"typeface": face}),
        ("a:buChar", {"char": char}),
    ):
        element = pPr.makeelement(qn(tag), attrs or {})
        if tag == "a:buClr":
            srgb = pPr.makeelement(qn("a:srgbClr"), {"val": f"{color:06X}"})
            element.append(srgb)
        pPr.append(element)


def _no_shadow(shape: Any) -> None:
    """No shadow, no theme fill, no inherited effects — the shape is exactly what we set.

    Two steps, because one isn't enough: an empty `<a:effectLst/>` states "no effects",
    and dropping the `<p:style>` element removes the reference to the theme's effect
    style that some renderers apply anyway. Without the second step the panels keep a
    faint drop shadow that says "Office autoshape" more loudly than anything else here.
    """
    from pptx.oxml.ns import qn

    spPr = shape._element.spPr
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)


def _rule(slide: Any, x: float, y: float, w: float, color: int, height: float = 0.045):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(color)
    bar.line.fill.background()
    _no_shadow(bar)
    return bar


def _panel(slide: Any, x: float, y: float, w: float, h: float, color: int):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(color)
    panel.line.fill.background()
    _no_shadow(panel)
    try:  # a softer corner than PowerPoint's default 16%
        panel.adjustments[0] = 0.04
    except (IndexError, AttributeError):
        pass
    return panel


def _place_title(
    slide: Any,
    text: str,
    theme: Theme,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    min_size: int = 18,
    line_spacing: Optional[float] = None,
) -> None:
    """Write the slide's title into the REAL title placeholder, moved and restyled.

    Free text boxes would give tidier code, but a slide whose title isn't a title
    placeholder has no title as far as PowerPoint's outline view, the slide sorter, or
    `read_presentation` are concerned — which is how a deck stops being revisable.
    """
    from pptx.util import Inches

    holder = slide.shapes.title
    if holder is None:
        frame = _box(slide, x, y, w, h)
    else:
        holder.left, holder.top = Inches(x), Inches(y)
        holder.width, holder.height = Inches(w), Inches(h)
        frame = holder.text_frame
        frame.word_wrap = True
        frame.clear()
    size = fitted_size(
        text, w, h, size, min_size=min_size, bold=True, line_spacing=line_spacing or 1.12
    )
    _write(
        frame,
        text,
        font=theme.heading,
        size=size,
        color=theme.ink,
        bold=True,
        first=True,
        line_spacing=line_spacing,
    )


def _slide_title(slide: Any, text: str, theme: Theme) -> None:
    _place_title(
        slide,
        text,
        theme,
        x=theme.margin_x,
        y=theme.title_top,
        w=theme.content_width,
        h=0.9,
        size=theme.size_slide_title,
    )
    _rule(slide, theme.margin_x, theme.title_top + 0.95, 1.15, theme.accent)


def _footer(slide: Any, theme: Theme, number: int) -> None:
    from pptx.enum.text import PP_ALIGN

    frame = _box(slide, 13.333 - theme.margin_x - 1.2, 6.95, 1.2, 0.3)
    # Named so `read_presentation` can skip it: a page number is chrome, and returning it
    # as a bullet would feed it back into the next revision as content.
    frame._parent.name = FOOTER_SHAPE_NAME
    para = _write(
        frame, str(number), font=theme.body, size=theme.size_footer, color=theme.muted, first=True
    )
    para.alignment = PP_ALIGN.RIGHT


# --------------------------------------------------------------------- layouts


def _paint_title(slide: Any, entry: dict, theme: Theme) -> None:
    _rule(slide, theme.margin_x, 2.35, 1.6, theme.accent, height=0.06)
    _place_title(
        slide,
        str(entry.get("title") or ""),
        theme,
        x=theme.margin_x,
        y=2.65,
        w=theme.content_width - 1.2,
        h=1.6,
        size=theme.size_deck_title,
        min_size=26,
        line_spacing=1.05,
    )
    subtitle = str(entry.get("subtitle") or "")
    if subtitle:
        sub = _box(slide, theme.margin_x, 4.45, theme.content_width - 1.2, 0.9)
        _write(sub, subtitle, font=theme.body, size=fitted_size(subtitle, theme.content_width - 1.2, 0.9, theme.size_body, min_size=13), color=theme.muted, first=True)


def _paint_section(slide: Any, entry: dict, theme: Theme, number: int) -> None:
    _panel(slide, 0, 2.55, 0.28, 1.9, theme.accent)
    _place_title(
        slide,
        str(entry.get("title") or ""),
        theme,
        x=theme.margin_x,
        y=2.6,
        w=theme.content_width,
        h=1.1,
        size=theme.size_section,
        min_size=22,
    )
    subtitle = str(entry.get("subtitle") or "")
    if subtitle:
        sub = _box(slide, theme.margin_x, 3.75, theme.content_width, 0.7)
        _write(sub, subtitle, font=theme.body, size=fitted_size(subtitle, theme.content_width, 0.7, theme.size_body, min_size=13), color=theme.muted, first=True)


def _bullets_height(bullets: list, width_in: float, base: int) -> float:
    """Estimated stack height (inches) of a bullet list whose top-level size is `base`."""
    total = 0.0
    for text, level in bullets:
        size = max(base - 3 * max(level, 0), 12)
        total += block_height(text, width_in, size, line_spacing=1.25)
        total += (14 if level == 0 else 8) / 72.0
    return total


def _fit_bullets(bullets: list, width_in: float, height_in: float, base: int) -> tuple[int, bool]:
    """(top-level size, fits) — shrink to 13pt before admitting defeat."""
    for candidate in range(base, 12, -1):
        if _bullets_height(bullets, width_in, candidate) <= height_in:
            return candidate, True
    return 13, False


def _paint_bullets(slide: Any, entry: dict, theme: Theme, bullets: list, warnings: list) -> None:
    height = theme.content_bottom - theme.content_top
    base, fits = _fit_bullets(bullets, theme.content_width, height, theme.size_body)
    if not fits:
        warnings.append(
            f"{len(bullets)} bullets don't fit even at minimum size — split this slide "
            "into two, or cut bullets."
        )
    frame = _box(slide, theme.margin_x, theme.content_top, theme.content_width, height)
    for i, (text, level) in enumerate(bullets):
        para = _write(
            frame,
            text,
            font=theme.body,
            size=max(base - 3 * max(level, 0), 12),
            color=theme.body_ink if level == 0 else theme.muted,
            first=(i == 0),
            space_after=14 if level == 0 else 8,
            line_spacing=1.25,
            level=level,
        )
        _bullet_marker(
            para,
            "—" if level == 0 else "·",
            theme.body,
            theme.accent if level == 0 else theme.muted,
            0.3 + 0.25 * level,
        )


def _paint_statement(slide: Any, entry: dict, theme: Theme) -> None:
    _rule(slide, theme.margin_x, 2.05, 1.6, theme.accent, height=0.06)
    _place_title(
        slide,
        str(entry.get("statement") or entry.get("title") or ""),
        theme,
        x=theme.margin_x,
        y=2.45,
        w=theme.content_width - 0.9,
        h=2.6,
        size=theme.size_statement,
        min_size=20,
        line_spacing=1.18,
    )
    support = str(entry.get("subtitle") or "")
    if support:
        sub = _box(slide, theme.margin_x, 5.2, theme.content_width - 0.9, 0.8)
        _write(
            sub,
            support,
            font=theme.body,
            size=fitted_size(support, theme.content_width - 0.9, 0.8, theme.size_body, min_size=13),
            color=theme.muted,
            first=True,
        )


def _paint_stat(slide: Any, entry: dict, theme: Theme, stats: list) -> None:
    from pptx.enum.text import MSO_ANCHOR

    count = max(1, min(len(stats), 3))
    gap = 0.45
    width = (theme.content_width - gap * (count - 1)) / count
    # The block is ~2.9in tall; centre it in the content region so the slide doesn't
    # read as top-heavy with a pool of white underneath.
    top = theme.content_top + max(0.0, ((theme.content_bottom - theme.content_top) - 2.9) / 2)
    for i, item in enumerate(stats[:count]):
        x = theme.margin_x + i * (width + gap)
        value = _box(slide, x, top, width, 1.3)
        value.vertical_anchor = MSO_ANCHOR.BOTTOM
        value_text = str(item.get("value") or "")
        # One line, always: a wrapped "-23pts" is worse than a smaller one.
        value_size = theme.size_stat
        while value_size > 32 and len(value_text) * _CHAR_FACTOR_BOLD * value_size / 72.0 > width:
            value_size -= 2
        _write(
            value,
            value_text,
            # Deliberately the BODY face, not the heading one: Georgia sets old-style
            # figures, so at 66pt "0pts" reads as the word "opts" and "-23pts" wobbles.
            # A stat slide exists to make one number unmistakable.
            font=theme.body,
            size=value_size,
            color=theme.accent,
            bold=True,
            first=True,
        )
        _rule(slide, x, top + 1.45, min(width, 1.0), theme.rule, height=0.02)
        label = _box(slide, x, top + 1.62, width, 1.3)
        _write(
            label,
            str(item.get("label") or ""),
            font=theme.body,
            size=theme.size_stat_label,
            color=theme.muted,
            first=True,
            line_spacing=1.3,
        )


def _paint_columns(
    slide: Any, entry: dict, theme: Theme, columns: list, warnings: list, *, carded: bool
) -> None:
    gap = 0.5
    count = max(1, min(len(columns), 2))
    width = (theme.content_width - gap * (count - 1)) / count
    top = theme.content_top
    available = theme.content_bottom - top
    pad_est = 0.32 if carded else 0.0
    heads = any(str(c.get("heading") or "") for c in columns[:count])
    text_room = available - 2 * pad_est - (0.62 if heads else 0)
    base = theme.size_body + (1 if carded else 0)
    fit = base
    for column in columns[:count]:
        pairs = _bullet_pairs(column.get("bullets") or [])
        size, fits = _fit_bullets(pairs, width - 2 * pad_est, text_room, base)
        fit = min(fit, size)
        if not fits:
            warnings.append(
                f"column '{column.get('heading') or '?'}' overflows even at minimum "
                "size — fewer or shorter bullets, or split the slide."
            )
    height = available
    if carded:
        # Size the cards to what's in them. A card that always runs to the bottom of the
        # slide leaves a pool of empty tint under three bullets and reads as unfinished.
        tallest = max(
            _bullets_height(_bullet_pairs(c.get("bullets") or []), width - 2 * pad_est, fit)
            for c in columns[:count]
        )
        height = min(available, 2 * pad_est + (0.62 if heads else 0) + tallest + 0.1)
        top += max(0.0, (available - height) / 2)
    for i, column in enumerate(columns[:count]):
        x = theme.margin_x + i * (width + gap)
        pad = 0.32 if carded else 0.0
        if carded:
            _panel(slide, x, top, width, height, theme.accent_soft if i == 0 else 0xF3F4F6)
        heading = str(column.get("heading") or "")
        y = top + pad
        if heading:
            head = _box(slide, x + pad, y, width - 2 * pad, 0.5)
            _write(
                head,
                heading,
                font=theme.heading,
                size=theme.size_body + 2,
                color=theme.accent if not carded else theme.ink,
                bold=True,
                first=True,
            )
            y += 0.62
        frame = _box(slide, x + pad, y, width - 2 * pad, height - (y - top) - pad)
        for j, (text, level) in enumerate(_bullet_pairs(column.get("bullets") or [])):
            para = _write(
                frame,
                text,
                font=theme.body,
                size=max(fit - 3 * max(level, 0), 12),
                color=theme.body_ink,
                first=(j == 0),
                space_after=10,
                line_spacing=1.25,
                level=level,
            )
            _bullet_marker(para, "—" if level == 0 else "·", theme.body, theme.accent, 0.26)


def _paint_quote(slide: Any, entry: dict, theme: Theme) -> None:
    words = str(entry.get("quote") or entry.get("title") or "")
    width = theme.content_width - 1.2
    size = fitted_size(words, width, 3.4, theme.size_quote, min_size=18, line_spacing=1.3)
    # The accent bar marks the quote, so it should be as tall as the quote is — a fixed
    # bar looms over a one-line quote and gets swamped by a four-line one.
    height = min(3.6, 0.15 + block_height(words, width, size, line_spacing=1.3))
    _panel(slide, theme.margin_x, 2.15, 0.09, height, theme.accent)
    frame = _box(slide, theme.margin_x + 0.45, 2.15, theme.content_width - 1.2, height + 0.1)
    _write(
        frame,
        "“" + words + "”",
        font=theme.heading,
        size=size,
        color=theme.ink,
        italic=True,
        first=True,
        line_spacing=1.3,
    )
    who = str(entry.get("attribution") or entry.get("subtitle") or "")
    if who:
        sub = _box(slide, theme.margin_x + 0.45, 2.15 + height + 0.25, theme.content_width - 1.2, 0.6)
        _write(sub, "— " + who, font=theme.body, size=theme.size_caption + 1, color=theme.muted, first=True)


def _paint_image(slide: Any, entry: dict, theme: Theme, picture_path: str) -> None:
    from pptx.util import Inches

    caption = str(entry.get("caption") or entry.get("takeaway") or "")
    box_top = theme.content_top
    box_height = (theme.content_bottom - box_top) - (0.6 if caption else 0.0)
    shape = slide.shapes.add_picture(picture_path, Inches(theme.margin_x), Inches(box_top))
    # Scale to fit the content box without distorting, then centre it.
    scale = min(
        Inches(theme.content_width) / shape.width, Inches(box_height) / shape.height, 1.0
    )
    shape.width, shape.height = int(shape.width * scale), int(shape.height * scale)
    shape.left = int(Inches(theme.margin_x) + (Inches(theme.content_width) - shape.width) / 2)
    shape.top = int(Inches(box_top) + (Inches(box_height) - shape.height) / 2)
    if caption:
        frame = _box(slide, theme.margin_x, theme.content_bottom - 0.45, theme.content_width, 0.5)
        _write(frame, caption, font=theme.body, size=theme.size_caption, color=theme.muted, first=True)


# ------------------------------------------------------------------ dispatch


def _bullet_pairs(items: Any) -> list[tuple[str, int]]:
    out: list[tuple[str, int]] = []
    for item in items or []:
        if isinstance(item, dict):
            out.append((str(item.get("text") or ""), int(item.get("level") or 0)))
        else:
            out.append((str(item), 0))
    return [(text, level) for text, level in out if text]


def paint(
    deck: Any,
    entry: dict,
    theme: Theme,
    *,
    number: int,
    resolve_image: Callable[[str], str],
) -> tuple[Any, list[str]]:
    """Draw one slide; returns (slide, layout warnings). A warning means the content
    doesn't fit its frame even at minimum size — a writing problem the caller (the
    model) must fix by splitting or cutting, because no renderer is available at
    runtime to screenshot-and-check."""
    warnings: list[str] = []
    kind = str(entry.get("layout") or "bullets").lower()
    if kind not in LAYOUTS:
        raise ValueError(
            f"unknown slide layout {kind!r}; expected one of: {', '.join(LAYOUTS)}"
        )

    title = str(entry.get("title") or "")
    # 'statement' puts its claim in the title placeholder, so it counts as titled even
    # when the words arrived in `statement` rather than `title`.
    titled = kind != "blank" and bool(title or (kind == "statement" and entry.get("statement")))
    chrome = kind not in ("title", "blank", "section")
    slide = (
        deck.slides.add_slide(_pick_layout(deck, "Title Only", "Title and Content"))
        if titled
        else _blank_slide(deck)
    )

    if chrome and title and kind != "statement":
        _slide_title(slide, title, theme)

    if kind == "title":
        _paint_title(slide, entry, theme)
    elif kind == "section":
        _paint_section(slide, entry, theme, number)
    elif kind == "bullets":
        _paint_bullets(slide, entry, theme, _bullet_pairs(entry.get("bullets")), warnings)
    elif kind == "statement":
        _paint_statement(slide, entry, theme)
    elif kind == "stat":
        stats = [s for s in (entry.get("stats") or []) if isinstance(s, dict)]
        if not stats:
            raise ValueError("a 'stat' slide needs stats: [{value, label}, …]")
        _paint_stat(slide, entry, theme, stats)
    elif kind in ("two_column", "comparison"):
        columns = [c for c in (entry.get("columns") or []) if isinstance(c, dict)]
        if len(columns) < 2:
            raise ValueError(
                f"a '{kind}' slide needs columns: [{{heading, bullets}}, {{heading, bullets}}]"
            )
        _paint_columns(slide, entry, theme, columns, warnings, carded=(kind == "comparison"))
    elif kind == "quote":
        _paint_quote(slide, entry, theme)
    elif kind == "image":
        image = str(entry.get("image") or "")
        if not image:
            raise ValueError("an 'image' slide needs image: <path to a png/jpg>")
        _paint_image(slide, entry, theme, resolve_image(image))

    if chrome:
        _footer(slide, theme, number)
    _drop_empty_placeholders(slide)
    return slide, warnings
