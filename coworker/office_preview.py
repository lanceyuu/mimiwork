"""Reading-quality previews of Word and PowerPoint files, and Word comments.

The GUI cannot render Office binaries itself, and until 2026-09-02 a .docx Mimi had
just written opened in Word or not at all. python-docx and python-pptx are already
in the sidecar for writing them, so they read them back too: headings, paragraphs,
lists, tables, pictures — the reading order, not the page layout. Every body
paragraph carries its index (``data-p``) and every slide its number
(``data-slide``), so a click in the preview can say exactly where a comment goes —
to Mimi as "paragraph 12, starting …", or into the file as a real Word comment.

Tracked changes and existing comments show as Word shows them (owner ask
2026-09-04): insertions underlined, deletions struck through, each with its author,
and a numbered marker where a comment sits. python-docx's ``paragraph.runs`` skips
every run inside ``w:ins``/``w:del``, so a reviewed document used to preview as if
all changes had been accepted — and the deletions were simply gone.
"""

from __future__ import annotations

import base64
import html
from pathlib import Path
from typing import Any

# Pictures are inlined as data URLs; past this much the preview would choke the
# JSON channel, so the remaining pictures become a placeholder.
MAX_INLINE_IMAGE_BYTES = 12 * 1024 * 1024


def _esc(text: str) -> str:
    return html.escape(text or "", quote=False)


class _Images:
    def __init__(self) -> None:
        self.used = 0

    def tag(self, blob: bytes, content_type: str, alt: str = "") -> str:
        if not blob:
            return ""
        if self.used + len(blob) > MAX_INLINE_IMAGE_BYTES:
            return '<span class="doc-img-skipped">[picture not shown — too large for the preview]</span>'
        self.used += len(blob)
        data = base64.b64encode(blob).decode("ascii")
        return f'<img src="data:{content_type or "image/png"};base64,{data}" alt="{_esc(alt)}">'


# -- Word ---------------------------------------------------------------------------


class _DocCtx:
    """What every paragraph needs while a document renders: the document (for image
    parts), the image budget, the comments by id, and the change tally for the legend."""

    def __init__(self, doc: Any) -> None:
        self.doc = doc
        self.images = _Images()
        self.comments: dict[str, tuple[str, str]] = {}
        self.ins = 0
        self.dels = 0
        self.comment_marks = 0
        try:
            for c in doc.comments:
                self.comments[str(c.comment_id)] = (c.author or "", c.text or "")
        except Exception:  # a file without a comments part, or an older python-docx
            pass


def _who(el: Any) -> str:
    from docx.oxml.ns import qn

    author = el.get(qn("w:author")) or "unknown"
    date = (el.get(qn("w:date")) or "")[:10]
    return f"{author}, {date}" if date else author


def _runs_html(el: Any, p: Any, ctx: _DocCtx) -> str:
    """The paragraph's runs in document order, tracked changes included."""
    from docx.oxml.ns import qn
    from docx.text.run import Run

    parts: list[str] = []
    for child in el.iterchildren():
        tag = child.tag
        if tag == qn("w:r"):
            parts.append(_run_html(Run(child, p), ctx.doc, ctx.images))
            ref = child.find(qn("w:commentReference"))
            if ref is not None:
                cid = str(ref.get(qn("w:id")))
                author, text = ctx.comments.get(cid, ("", ""))
                ctx.comment_marks += 1
                label = f"{author}: {text}" if author else text
                parts.append(
                    f'<sup class="doc-comment" data-comment="{_esc(cid)}" title="{_esc(label)}">'
                    f"{ctx.comment_marks}</sup>"
                )
        elif tag in (qn("w:ins"), qn("w:moveTo")):
            ctx.ins += 1
            verb = "Inserted" if tag == qn("w:ins") else "Moved here"
            parts.append(
                f'<ins class="doc-ins" title="{verb} by {_esc(_who(child))}">{_runs_html(child, p, ctx)}</ins>'
            )
        elif tag in (qn("w:del"), qn("w:moveFrom")):
            ctx.dels += 1
            verb = "Deleted" if tag == qn("w:del") else "Moved away"
            text = "".join(t.text or "" for t in child.iter(qn("w:delText")))
            parts.append(f'<del class="doc-del" title="{verb} by {_esc(_who(child))}">{_esc(text)}</del>')
        elif tag in (qn("w:hyperlink"), qn("w:smartTag"), qn("w:sdt"), qn("w:sdtContent")):
            parts.append(_runs_html(child, p, ctx))
    return "".join(parts)


def _run_html(run: Any, doc: Any, images: _Images) -> str:
    from docx.oxml.ns import qn

    parts: list[str] = []
    for blip in run._r.iter(qn("a:blip")):
        rid = blip.get(qn("r:embed"))
        part = doc.part.related_parts.get(rid) if rid else None
        if part is not None:
            parts.append(images.tag(part.blob, getattr(part, "content_type", "")))
    text = _esc(run.text)
    if text:
        if run.bold:
            text = f"<b>{text}</b>"
        if run.italic:
            text = f"<i>{text}</i>"
        if run.underline:
            text = f"<u>{text}</u>"
        parts.append(text)
    return "".join(parts)


def _paragraph_html(p: Any, ctx: _DocCtx, index: int | None) -> tuple[str, str]:
    """(kind, html) — kind is 'h', 'li' or 'p', so consecutive list items can be grouped."""
    from docx.oxml.ns import qn

    style = (p.style.name if p.style is not None else "") or ""
    inner = _runs_html(p._p, p, ctx) or _esc(p.text)
    attr = f' data-p="{index}"' if index is not None else ""
    # A deleted paragraph mark: the whole paragraph is going away.
    ppr = p._p.pPr
    mark = ppr.find(qn("w:rPr")) if ppr is not None else None
    if mark is not None and mark.find(qn("w:del")) is not None:
        attr += ' class="doc-del-p"'
    if style.startswith("Heading"):
        level = "".join(ch for ch in style if ch.isdigit()) or "1"
        n = min(4, max(1, int(level)))
        return "h", f"<h{n}{attr}>{inner}</h{n}>"
    if style == "Title":
        return "h", f"<h1{attr}>{inner}</h1>"
    numbered = p._p.pPr is not None and p._p.pPr.numPr is not None
    if numbered or "List" in style:
        return "li", f"<li{attr}>{inner}</li>"
    if not inner.strip():
        return "p", ""
    return "p", f"<p{attr}>{inner}</p>"


def _table_html(table: Any, ctx: _DocCtx) -> str:
    rows: list[str] = []
    for r, row in enumerate(table.rows):
        cells: list[str] = []
        for cell in row.cells:
            # Cell text stays inline (one line per paragraph): a <p> per cell paragraph
            # made every table twice as tall as Word shows it.
            lines = [_runs_html(p._p, p, ctx) or _esc(p.text) for p in cell.paragraphs]
            inner = "<br>".join(line for line in lines if line.strip()) or _esc(cell.text)
            tag = "th" if r == 0 else "td"
            cells.append(f"<{tag}>{inner}</{tag}>")
        rows.append("<tr>" + "".join(cells) + "</tr>")
    return "<table>" + "".join(rows) + "</table>"


def docx_to_html(path: str | Path) -> dict[str, Any]:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(str(path))
    ctx = _DocCtx(doc)
    out: list[str] = []
    in_list = False
    index = 0
    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            kind, h = _paragraph_html(Paragraph(child, doc), ctx, index)
            index += 1
        elif child.tag == qn("w:tbl"):
            kind, h = "p", _table_html(Table(child, doc), ctx)
        else:
            continue
        if kind == "li" and not in_list:
            out.append("<ul>")
            in_list = True
        elif kind != "li" and in_list:
            out.append("</ul>")
            in_list = False
        if h:
            out.append(h)
    if in_list:
        out.append("</ul>")
    changes = {"insertions": ctx.ins, "deletions": ctx.dels, "comments": ctx.comment_marks}
    if ctx.ins or ctx.dels or ctx.comment_marks:
        bits = []
        if ctx.ins or ctx.dels:
            bits.append(
                "Track changes: "
                + ", ".join(
                    f"{n} {word}{'' if n == 1 else 's'}"
                    for n, word in ((ctx.ins, "insertion"), (ctx.dels, "deletion"))
                    if n
                )
            )
        if ctx.comment_marks:
            bits.append(f"{ctx.comment_marks} comment{'' if ctx.comment_marks == 1 else 's'}")
        legend = " · ".join(bits)
        out.insert(
            0,
            '<div class="doc-changes"><span class="doc-changes-ins">inserted</span>'
            '<span class="doc-changes-del">deleted</span>' + _esc(legend) + "</div>",
        )
    return {"html": "\n".join(out), "paragraphs": index, "changes": changes}


def add_word_comment(
    path: str | Path, paragraph: int, text: str, author: str = "MimiWork"
) -> dict[str, Any]:
    """A real Word comment on body paragraph `paragraph` (0-based), saved in place."""
    from docx import Document

    doc = Document(str(path))
    paras = doc.paragraphs
    if not 0 <= paragraph < len(paras):
        raise ValueError(f"no paragraph {paragraph + 1} — the document has {len(paras)}")
    p = paras[paragraph]
    runs = list(p.runs) or [p.add_run("")]
    initials = "".join(w[0] for w in author.split()[:2]).upper() or "M"
    doc.add_comment(runs, text=text, author=author, initials=initials)
    doc.save(str(path))
    return {"paragraph": paragraph, "comments": len(list(doc.comments))}


# -- PowerPoint ---------------------------------------------------------------------


def pptx_to_html(path: str | Path) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    images = _Images()
    out: list[str] = []
    n = 0
    for n, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        parts: list[str] = [f'<section class="slide" data-slide="{n}"><div class="slide-no">Slide {n}</div>']
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    parts.append(images.tag(shape.image.blob, shape.image.content_type))
                except Exception:
                    continue
                continue
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = []
                for r, row in enumerate(shape.table.rows):
                    tag = "th" if r == 0 else "td"
                    rows.append("<tr>" + "".join(f"<{tag}>{_esc(c.text)}</{tag}>" for c in row.cells) + "</tr>")
                parts.append("<table>" + "".join(rows) + "</table>")
                continue
            if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
                continue
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                parts.append(f"<h2>{_esc(shape.text_frame.text)}</h2>")
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs) or para.text
                if not text.strip():
                    continue
                indent = f' style="margin-left:{para.level * 16}px"' if para.level else ""
                parts.append(f"<p{indent}>{_esc(text)}</p>")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f'<aside class="slide-notes">Notes: {_esc(notes)}</aside>')
        parts.append("</section>")
        out.append("".join(parts))
    return {"html": "\n".join(out), "slides": n}
